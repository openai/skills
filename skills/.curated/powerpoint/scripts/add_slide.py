#!/usr/bin/env python3
"""
Add a slide to an existing PowerPoint presentation.

Usage:
    python add_slide.py presentation.pptx --title "New Slide" --bullets "Point 1" "Point 2"
    python add_slide.py presentation.pptx --title "Title Only" --layout 5
    python add_slide.py presentation.pptx --image photo.png --position 2
    
    Or import and use programmatically:
    >>> from add_slide import add_slide
    >>> add_slide('presentation.pptx', title='New Slide', bullets=['Point 1', 'Point 2'])
"""

import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    exit(1)


def add_slide(
    pptx_path: str,
    title: str = "",
    bullets: list[str] | None = None,
    image_path: str | None = None,
    layout: int = 1,
    position: int | None = None,
    output_path: str | None = None
) -> str:
    """
    Add a slide to a PowerPoint presentation.
    
    Args:
        pptx_path: Path to the .pptx file
        title: Slide title
        bullets: List of bullet points
        image_path: Path to image to add
        layout: Slide layout index (0=title, 1=bullets, 5=title only, 6=blank)
        position: Insert position (1-based), None for end
        output_path: Save path (default: overwrite original)
    
    Returns:
        Path to the saved file
    """
    prs = Presentation(pptx_path)
    
    # Add new slide
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    
    # Set title
    if title and slide.shapes.title:
        slide.shapes.title.text = title
    
    # Add bullets
    if bullets and layout == 1:
        tf = slide.placeholders[1].text_frame
        tf.text = bullets[0] if bullets else ""
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    
    # Add image
    if image_path:
        left = Inches(1)
        top = Inches(2) if title else Inches(1)
        slide.shapes.add_picture(image_path, left, top, height=Inches(4))
    
    # Move slide to position if specified
    if position is not None and position > 0:
        # Calculate the target position (0-indexed)
        target_idx = min(position - 1, len(prs.slides) - 1)
        current_idx = len(prs.slides) - 1
        
        # Move by manipulating the slide ID list
        slide_id = prs.slides._sldIdLst[-1]
        prs.slides._sldIdLst.remove(slide_id)
        prs.slides._sldIdLst.insert(target_idx, slide_id)
    
    save_path = output_path or pptx_path
    prs.save(save_path)
    return str(Path(save_path).resolve())


def main():
    parser = argparse.ArgumentParser(description='Add slide to PowerPoint')
    parser.add_argument('pptx', help='Input .pptx file')
    parser.add_argument('--title', '-t', default='', help='Slide title')
    parser.add_argument('--bullets', '-b', nargs='+', help='Bullet points')
    parser.add_argument('--image', '-i', help='Image to add')
    parser.add_argument('--layout', '-l', type=int, default=1,
                        help='Layout: 0=title, 1=bullets, 5=title only, 6=blank')
    parser.add_argument('--position', '-p', type=int, help='Insert position (1-based)')
    parser.add_argument('--output', '-o', help='Output file (default: overwrite)')
    
    args = parser.parse_args()
    
    result = add_slide(
        args.pptx,
        title=args.title,
        bullets=args.bullets,
        image_path=args.image,
        layout=args.layout,
        position=args.position,
        output_path=args.output
    )
    print(f"Updated: {result}")


if __name__ == '__main__':
    main()
