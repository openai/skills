#!/usr/bin/env python3
"""
Create a PowerPoint presentation from structured data.

Usage:
    python create_presentation.py output.pptx --title "My Title" --slides slides.json
    
    Or import and use programmatically:
    >>> from create_presentation import create_presentation
    >>> create_presentation('output.pptx', title='My Presentation', slides=slides_data)
"""

import argparse
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    exit(1)


def create_presentation(
    output_path: str,
    title: str = "Presentation",
    subtitle: str = "",
    slides: list[dict] | None = None
) -> str:
    """
    Create a PowerPoint presentation.
    
    Args:
        output_path: Path to save the .pptx file
        title: Title for the title slide
        subtitle: Subtitle for the title slide
        slides: List of slide dicts with keys:
            - title (str): Slide title
            - content (str | list): Text content or list of bullet points
            - layout (int, optional): Slide layout index (default: 1 for bullets)
    
    Returns:
        Path to the created file
    """
    prs = Presentation()
    
    # Add title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if subtitle:
        title_slide.placeholders[1].text = subtitle
    
    # Add content slides
    if slides:
        for slide_data in slides:
            layout_idx = slide_data.get('layout', 1)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            
            # Set title
            if 'title' in slide_data and slide.shapes.title:
                slide.shapes.title.text = slide_data['title']
            
            # Set content
            content = slide_data.get('content')
            if content and layout_idx == 1:
                tf = slide.placeholders[1].text_frame
                if isinstance(content, list):
                    tf.text = content[0] if content else ""
                    for item in content[1:]:
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = 0
                else:
                    tf.text = str(content)
    
    prs.save(output_path)
    return str(Path(output_path).resolve())


def main():
    parser = argparse.ArgumentParser(description='Create a PowerPoint presentation')
    parser.add_argument('output', help='Output .pptx file path')
    parser.add_argument('--title', default='Presentation', help='Title slide title')
    parser.add_argument('--subtitle', default='', help='Title slide subtitle')
    parser.add_argument('--slides', help='JSON file with slides data')
    
    args = parser.parse_args()
    
    slides = None
    if args.slides:
        with open(args.slides) as f:
            slides = json.load(f)
    
    result = create_presentation(
        args.output,
        title=args.title,
        subtitle=args.subtitle,
        slides=slides
    )
    print(f"Created: {result}")


if __name__ == '__main__':
    main()
