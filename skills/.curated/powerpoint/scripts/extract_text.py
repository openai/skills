#!/usr/bin/env python3
"""
Extract all text content from a PowerPoint presentation.

Usage:
    python extract_text.py presentation.pptx
    python extract_text.py presentation.pptx --output text.txt
    python extract_text.py presentation.pptx --json
    
    Or import and use programmatically:
    >>> from extract_text import extract_text
    >>> text = extract_text('presentation.pptx')
"""

import argparse
import json
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    exit(1)


def extract_text(pptx_path: str, include_notes: bool = True) -> list[dict]:
    """
    Extract all text from a PowerPoint presentation.
    
    Args:
        pptx_path: Path to the .pptx file
        include_notes: Whether to include speaker notes
    
    Returns:
        List of dicts with slide number, title, content, and notes
    """
    prs = Presentation(pptx_path)
    slides_data = []
    
    for idx, slide in enumerate(prs.slides, 1):
        slide_info = {
            'slide_number': idx,
            'title': '',
            'content': [],
            'notes': ''
        }
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = '\n'.join(
                    para.text for para in shape.text_frame.paragraphs
                    if para.text.strip()
                )
                if text:
                    # Check if this is the title shape
                    if shape == slide.shapes.title:
                        slide_info['title'] = text
                    else:
                        slide_info['content'].append(text)
        
        # Extract notes
        if include_notes and slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                slide_info['notes'] = notes_frame.text
        
        slides_data.append(slide_info)
    
    return slides_data


def format_as_text(slides_data: list[dict]) -> str:
    """Format extracted data as plain text."""
    lines = []
    for slide in slides_data:
        lines.append(f"--- Slide {slide['slide_number']} ---")
        if slide['title']:
            lines.append(f"Title: {slide['title']}")
        for content in slide['content']:
            lines.append(content)
        if slide['notes']:
            lines.append(f"\n[Notes: {slide['notes']}]")
        lines.append("")
    return '\n'.join(lines)


def main():
    parser = argparse.ArgumentParser(description='Extract text from PowerPoint')
    parser.add_argument('pptx', help='Input .pptx file')
    parser.add_argument('--output', '-o', help='Output file (default: stdout)')
    parser.add_argument('--json', action='store_true', help='Output as JSON')
    parser.add_argument('--no-notes', action='store_true', help='Exclude speaker notes')
    
    args = parser.parse_args()
    
    slides_data = extract_text(args.pptx, include_notes=not args.no_notes)
    
    if args.json:
        output = json.dumps(slides_data, indent=2)
    else:
        output = format_as_text(slides_data)
    
    if args.output:
        Path(args.output).write_text(output)
        print(f"Saved to: {args.output}")
    else:
        print(output)


if __name__ == '__main__':
    main()
